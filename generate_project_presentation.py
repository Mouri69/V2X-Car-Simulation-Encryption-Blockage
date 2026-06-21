from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "V2X Communication Simulation with MITM Attack & Post-Quantum Security"
    subtitle.text = "Graduation Project\nYour Name Here\nUniversity Name Here"

    # --- Slide 2: Project Overview ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Project Overview"
    tf = body_shape.text_frame
    tf.text = "What is this simulation?"
    p = tf.add_paragraph()
    p.text = "• SUMO-based (Simulation of Urban MObility) traffic simulation"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Models Vehicle-to-Vehicle (V2V) communication"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Shows both insecure classical encryption and secure post-quantum encryption"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Features a Man-in-the-Middle (MITM) attacker vehicle"
    p.level = 1

    # --- Slide 3: Original Simulation ---
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Original Simulation (Starting Point)"
    tf = body_shape.text_frame
    tf.text = "Features before your modifications:"
    p = tf.add_paragraph()
    p.text = "• Traffic simulation with 12 vehicles"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• MITM attacker that can intercept messages"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Switch from classical to post-quantum crypto at 60s"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Basic vehicle info (position, speed)"
    p.level = 1

    # --- Slide 4: Your Modifications Part 1 ---
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Your Modifications (Part 1)"
    tf = body_shape.text_frame
    tf.text = "1. Changed crypto switch from 60s → 90s (more attack demo time)"
    p = tf.add_paragraph()
    p.text = "2. Added vehicle owner names (James, Don, Sarah, etc.)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "3. Added sensitive personal/vehicle info:"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "   • Car Serial Numbers (VINs)"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "   • Owner Security Numbers (simulated SSNs)"
    p.level = 2

    # --- Slide 5: Your Modifications Part 2 ---
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Your Modifications (Part 2)"
    tf = body_shape.text_frame
    tf.text = "4. Swapped button functionality for clearer workflow:"
    p = tf.add_paragraph()
    p.text = "   • Button 4: Get ALL ENCOUNTERED Vehicles' Data"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "   • Button 5: Get SELECTED TARGET Data Only"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "5. Updated target selection UI to show [ENCOUNTERED]"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "6. Fixed Tkinter thread-safety issue"
    p.level = 1

    # --- Slide 6: Your Modifications Part 3 ---
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Your Modifications (Part 3)"
    tf = body_shape.text_frame
    tf.text = "7. Enhanced post-quantum security enforcement:"
    p = tf.add_paragraph()
    p.text = "   • Blocks all attacks AND info retrieval"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "   • Releases controlled vehicles back to SUMO"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "   • Logs clear [SAFE]/[ACCESS DENIED] messages"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "8. Added collision handling"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "9. Slowed simulation for better visibility"
    p.level = 1

    # --- Slide 7: Demo Flow ---
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Simulation Demo Flow"
    tf = body_shape.text_frame
    tf.text = "1. 0-90s (Classical Encryption):"
    p = tf.add_paragraph()
    p.text = "   • Purple Car encounters vehicles"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "   • Perform attacks & retrieve data"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "2. 90s (Crypto Switch):"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "   • All attacks/data retrieval BLOCKED"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "3. Post-90s (Post-Quantum Encryption):"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "   • Controlled vehicles released"
    p.level = 1

    # --- Slide 8: Key Results ---
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Key Results & Conclusion"
    tf = body_shape.text_frame
    tf.text = "Before 90s (Vulnerable):"
    p = tf.add_paragraph()
    p.text = "• Attacks succeed"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Sensitive info retrievable"
    p.level = 1
    p = tf.add_paragraph()
    p.text = ""
    p.level = 0
    p = tf.add_paragraph()
    p.text = "After 90s (Protected):"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• All attacks/retrieval FAIL"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Your Contribution: Turned basic sim into realistic, educational V2X security demo!"
    p.level = 0

    # --- Slide 9: Future Work ---
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Future Work (Optional Extra Credit)"
    tf = body_shape.text_frame
    tf.text = "• Add real post-quantum algorithms (CRYSTALS-Kyber)"
    p = tf.add_paragraph()
    p.text = "• Add more attack types (message spoofing, jamming)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Add multiple attackers or defensive vehicles"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Integrate with ns-3 for real-time network simulation"
    p.level = 1

    # Save the presentation
    prs.save('V2X_Post_Quantum_Security_Presentation.pptx')
    print("Presentation generated: V2X_Post_Quantum_Security_Presentation.pptx")
    print("\nImage Prompts for Slides:")
    print("Slide 1: Image of connected vehicles on a road")
    print("Slide 3: Screenshot of your simulation's initial state")
    print("Slide 4: Screenshot of vehicle data showing VINs and SSNs")
    print("Slide 6: Screenshot of crypto switch message saying 'ACCESS DENIED'")
    print("Slide 7: Screenshot of attacker GUI with successful attack")
    print("Slide 8: Side-by-side screenshots of attack success vs failure")

if __name__ == "__main__":
    create_presentation()