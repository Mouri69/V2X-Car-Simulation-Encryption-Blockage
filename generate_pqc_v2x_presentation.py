from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

slides = [
    {
        "title": "Post‑Quantum Cryptography for Secure V2X",
        "bullets": [
            "Graduation project: PQC for V2X communication",
            "Research paper and SUMO/Python simulation",
            "Focus: safety, longevity, quantum‑resistant security"
        ],
        "visual": "Title slide with V2X diagram (cars, RSU, cloud)"
    },
    {
        "title": "V2X Security Requirements",
        "bullets": [
            "Low latency and high reliability",
            "Integrity and authenticity for safety messages",
            "Confidentiality and privacy for data exchange"
        ],
        "visual": "Standards map (BSM/DENM) and message flow diagram"
    },
    {
        "title": "Quantum Threat Model",
        "bullets": [
            "Shor breaks RSA/ECC via factoring/discrete logs",
            "Grover weakens brute‑force for symmetric keys",
            "Harvest‑now, decrypt‑later risk for long‑lived systems"
        ],
        "visual": "Algorithm impact matrix (Shor/Grover vs RSA/ECC/AES)"
    },
    {
        "title": "Life‑or‑Death Risks in V2X",
        "bullets": [
            "Forged emergency brake can cause collisions",
            "Fake accident alerts misroute vehicles",
            "Safety requires robust authenticity and trust"
        ],
        "visual": "Incident sequence with forged message leading to pile‑up"
    },
    {
        "title": "PQC and NIST Standardization",
        "bullets": [
            "Lattice‑based Module‑LWE/Module‑SIS primitives",
            "NIST PQC: Kyber (KEM), Dilithium (signature)",
            "Goal: resist known quantum attacks"
        ],
        "visual": "NIST PQC timeline and algorithm families"
    },
    {
        "title": "CRYSTALS‑Kyber (Key Exchange)",
        "bullets": [
            "Efficient KEM for session key establishment",
            "Balanced performance and security",
            "Standardized and widely supported"
        ],
        "visual": "KEM handshake diagram between vehicles"
    },
    {
        "title": "CRYSTALS‑Dilithium (Signatures)",
        "bullets": [
            "Strong authenticity with moderate compute",
            "Suitable for millisecond decisions",
            "Robust against quantum cryptanalysis"
        ],
        "visual": "Signed BSM flow and verification path"
    },
    {
        "title": "Proposed V2X Crypto Architecture",
        "bullets": [
            "Kyber for key encapsulation and session keys",
            "AES‑GCM payload after PQ key agreement",
            "Dilithium for message integrity/authenticity"
        ],
        "visual": "Layered stack: PQC + symmetric encryption"
    },
    {
        "title": "Performance Trade‑offs",
        "bullets": [
            "Packet size grows to ~1–3 KB vs ~32 B",
            "Higher ECU load for lattice math",
            "Network throughput and congestion planning"
        ],
        "visual": "Bar chart of packet size and CPU cost vs classical"
    },
    {
        "title": "Simulation Setup (SUMO/Python)",
        "bullets": [
            "Map, routes, Red/Green/Blue/Purple vehicles",
            "Purple attacker with MITM capabilities",
            "Mode schedule: classical then PQC"
        ],
        "visual": "SUMO screenshot with labeled vehicles"
    },
    {
        "title": "Classical Mode: Attack Demonstration",
        "bullets": [
            "MITM intercepts and forges safety messages",
            "Shor‑simulated break of RSA/ECC key exchange",
            "Logs confirm forged and intercepted status"
        ],
        "visual": "Timeline of intercepted and forged communications"
    },
    {
        "title": "PQC Mode: Defense",
        "bullets": [
            "Kyber/Dilithium resist Shor/Grover attempts",
            "MITM blocked; authenticity preserved",
            "Advisory replaces attack banner for normal traffic"
        ],
        "visual": "Comparison table: classical vs PQC outcomes"
    },
    {
        "title": "Migration Strategy for OEMs",
        "bullets": [
            "Hybrid stack: PQ KEM + AES‑256 payload",
            "Roadmap for ECU hardware acceleration",
            "Standards and certification alignment"
        ],
        "visual": "Roadmap graphic (short/mid/long‑term milestones)"
    },
    {
        "title": "Conclusions and Future Work",
        "bullets": [
            "PQC is necessary for safety‑critical V2X",
            "Simulation validates attacker vs defense",
            "Next: real PQC libs and hardware tests"
        ],
        "visual": "Summary with key takeaways"
    },
    {
        "title": "Demo vs Real Implementation Checks",
        "bullets": [
            "Are real RSA/ECC keys generated via library APIs",
            "Is ML‑KEM/Dilithium used for lattice math",
            "Is attack success tied to actual crypto outcomes"
        ],
        "visual": "Checklist slide for code verification"
    }
]

def add_slide(prs, title, bullets, visual):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
        else:
            p = body.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(20)
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1))
    tf = tx.text_frame
    tf.text = "Visual: " + visual
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT

def main():
    prs = Presentation()
    for s in slides:
        add_slide(prs, s["title"], s["bullets"], s["visual"])
    prs.save("PQC_V2X_Presentation.pptx")

if __name__ == "__main__":
    main()
